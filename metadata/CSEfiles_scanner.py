#!/usr/bin/env python
# coding: utf-8

import os, json
import re
import docx2txt
import fitz
import zipfile
import PyPDF2
import minecart
import base64
import shutil
import xlrd
import csv
import win32com.client
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from StringIO import StringIO
except ImportError:
    from io import StringIO

def readDOCX(file_name):
    try:
        if file_name.endswith('.docx'):
            return docx2txt.process(file_name)
        elif file_name.endswith('.doc'):
            word = win32com.client.Dispatch("Word.Application")
            word.visible = False
            wb = word.Documents.Open(file_name)
            doc = word.ActiveDocument.Range().Text
            wb.Close(True)
            return doc
    except FileNotFoundError:
        print("Wrong file or file path")
    else:
        return null


def readPDF(file_name):
    try:
        # creating an object 
        pdf = fitz.open(file_name)
        content = ''
        for page_index in range(pdf.pageCount):
            text = json.loads(pdf.getPageText(page_index, output='json'))
            for block in text['blocks']:
                if 'lines' not in block:
                    # Skip blocks without text
                    continue
                for line in block['lines']:
                    for span in line['spans']:
                        content = content + str(span['text'].encode('utf-8'))[2:-1]
        content = re.sub(r"\r\n", " ", content)
        pdf.close()
        return content
    except FileNotFoundError:
        print("Wrong file or file path")
    else:
        return null

def readEXCEL(file_name):
    try:
        workbook = xlrd.open_workbook(file_name, on_demand=True)
        sheet_content = []
        sheet_names = workbook.sheet_names()
        for sheet_no in range(0, len(sheet_names)):
            for rownum in range(
                    workbook.sheet_by_index(sheet_no).nrows):  # sh1.nrows -> number of rows (ncols -> num columns)
                sheet_content = sheet_content + workbook.sheet_by_index(sheet_no).row_values(rownum)
        return ' '.join(str(sc) for sc in sheet_content)
    except FileNotFoundError:
        print("Wrong file or file path")
    else:
        return null

def readCSV(file_name):
    try:
        with open(file_name, 'r') as csvfile:
            spamreader = csv.reader(csvfile, delimiter=' ', quotechar='|')
            tags = []
            for row in spamreader:
                tags = tags + row[0].split(',')
            csvfile.close()
        return " ".join(tags)
    except FileNotFoundError:
        print("Wrong file or file path")
    else:
        return null

def readTXT(file_name):
    try:
        f = open(file_name, "r")
        return f.read()
    except FileNotFoundError:
        print("Wrong file or file path")
    else:
        return null

def readPPTX(file_name):
    try:
        prs = Presentation(file_name)
        text_runs = ''

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs += run.text + ' '

        return text_runs
    except FileNotFoundError:
        return "File Not Found"

'''
The following functions are for extracting images from different document files
'''
def getFilteredPath(file_name, ext=True):
    symbols = (os.sep, '\/', ':', '*', '?', '\"', '<', '>', '|', '.', ' ', '-')
    for symbol in symbols:
        file_name = file_name.replace(symbol, '_')
    if ext:
        return file_name + r".json"
    else:
        return file_name

def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape

def extractPPTXImages(file_name, output_folder, img_dir):
    try:
        output_folder = os.path.join(img_dir, getFilteredPath(file_name, False))
        if (os.path.exists(output_folder)):
            shutil.rmtree(output_folder)
        os.mkdir(output_folder)
        imglist = []
        i = 1
        for picture in iter_picture_shapes(Presentation(file_name)):
            image = picture.image
            image_bytes = image.blob
            with open(os.path.join(output_folder, 'image' + str(i) + '.png'), 'wb') as f:
                f.write(image_bytes)
            imglist.append(os.path.join(output_folder, 'image' + str(i) + '.png'))
            i = i + 1
        
        return imglist
    except FileNotFoundError:
        return "File Not Found"

def extractPDFImages(file_name, output_folder, img_dir):
    
    try:
        pdffile = open(file_name, 'rb')
        fileHandler = PyPDF2.PdfFileReader(open(file_name, "rb"))
        # print(fileHandler.numPages)
        doc = minecart.Document(pdffile)
        imglist = []
        output_folder = os.path.join(img_dir, getFilteredPath(file_name, False))

        if (os.path.exists(output_folder)):
            shutil.rmtree(output_folder)
        os.mkdir(output_folder)

        j = 1
        for i in range(0, fileHandler.numPages):
            try:
                page = doc.get_page(i)
                # print(page.images[0])
                # print(len(page.images))
                for image in page.images:
                    byteArray = image.obj.get_data()
                    with open(os.path.join(output_folder, 'image' + str(j) + '.png'), 'wb') as f:
                        f.write(byteArray)
                    imglist.append(os.path.join(output_folder, 'image' + str(j) + '.png'))
                    j = j + 1
            except:
                continue
        pdffile.close()
        #print(imglist)
        return imglist
    except Exception as e:
        print(str(e))

def extractDOCXandPPTXImages(file_name, output_folder, img_dir):
    try:
        temp_folder = os.path.join(output_folder, 'temp')
        if (os.path.exists(temp_folder)):
            shutil.rmtree(temp_folder)
        os.mkdir(temp_folder)
        efile = zipfile.ZipFile(file_name)
        efile.extractall(temp_folder)
        EmbeddedFiles = zipfile.ZipFile(file_name).namelist()
        # print(EmbeddedFiles)
        imglist = []

        if (os.path.exists(os.path.join(img_dir, getFilteredPath(file_name, False)))):
            shutil.rmtree(os.path.join(img_dir, getFilteredPath(file_name, False)))
        os.mkdir(os.path.join(img_dir, getFilteredPath(file_name, False)))

        for F in EmbeddedFiles:
            if '.jpg' in F or '.jpeg' in F or '.png' in F or '.m4v' in F:
                F = F.replace('/', '\\')
                # print(os.path.join(temp_folder,F))
                shutil.copy(os.path.join(temp_folder, F), os.path.join(img_dir, getFilteredPath(file_name, False)))
                imglist.append(os.path.join(img_dir, getFilteredPath(file_name, False)) + os.sep + os.path.basename(F))
        shutil.rmtree(temp_folder)
        return imglist
    except FileNotFoundError:
        return None

def extractXLS(file_name, output_folder, img_dir):
    return ['None']
def extractXLSX(file_name, output_folder, img_dir):
    return ['None']